"""
文档解析服务
支持 PDF、DOCX、XLSX、PPTX 文件解析为文本
基于 Skills 目录中的工具和库
"""
import logging
from pathlib import Path
from typing import Optional, List, Dict, Any
import tempfile
import re

logger = logging.getLogger(__name__)


class DocumentParser:
    """文档解析器，支持多种Office文件格式"""
    
    @staticmethod
    def parse_pdf(file_path: str) -> str:
        """
        解析PDF文件提取文本
        使用 pdfplumber 库（Skills/pdf/SKILL.md）
        """
        try:
            import pdfplumber
            
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {i+1} ---\n{page_text}")
            
            if not text_parts:
                raise ValueError("PDF文件中未提取到任何文本")
            
            return "\n\n".join(text_parts)
            
        except ImportError:
            raise RuntimeError("缺少依赖: pip install pdfplumber")
        except Exception as e:
            logger.error(f"PDF解析失败: {e}")
            raise ValueError(f"PDF解析失败: {str(e)}")
    
    @staticmethod
    def parse_docx(file_path: str) -> str:
        """
        解析DOCX文件提取文本
        使用 python-docx 库（Skills/docx/SKILL.md）
        """
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_parts = []
            
            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            if not text_parts:
                raise ValueError("DOCX文件中未提取到任何文本")
            
            return "\n\n".join(text_parts)
            
        except ImportError:
            raise RuntimeError("缺少依赖: pip install python-docx")
        except Exception as e:
            logger.error(f"DOCX解析失败: {e}")
            raise ValueError(f"DOCX解析失败: {str(e)}")
    
    @staticmethod
    def parse_xlsx(file_path: str) -> str:
        """
        解析XLSX文件提取文本
        使用 openpyxl 库（Skills/xlsx/SKILL.md）
        """
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    # 过滤空行
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(val.strip() for val in row_values):
                        text_parts.append(" | ".join(row_values))
            
            if len(text_parts) <= len(wb.sheetnames):
                raise ValueError("XLSX文件中未提取到任何数据")
            
            return "\n".join(text_parts)
            
        except ImportError:
            raise RuntimeError("缺少依赖: pip install openpyxl")
        except Exception as e:
            logger.error(f"XLSX解析失败: {e}")
            raise ValueError(f"XLSX解析失败: {str(e)}")
    
    @staticmethod
    def parse_pptx(file_path: str) -> str:
        """
        解析PPTX文件提取文本
        使用 python-pptx 库（Skills/pptx/SKILL.md）
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            
            for i, slide in enumerate(prs.slides):
                text_parts.append(f"--- Slide {i+1} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
                
                # 提取备注
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        text_parts.append(f"Notes: {notes}")
            
            if len(text_parts) <= len(prs.slides):
                raise ValueError("PPTX文件中未提取到任何文本")
            
            return "\n\n".join(text_parts)
            
        except ImportError:
            raise RuntimeError("缺少依赖: pip install python-pptx")
        except Exception as e:
            logger.error(f"PPTX解析失败: {e}")
            raise ValueError(f"PPTX解析失败: {str(e)}")
    
    @classmethod
    def parse_file(cls, file_path: str, file_type: Optional[str] = None) -> str:
        """
        根据文件类型自动选择解析器
        
        Args:
            file_path: 文件路径
            file_type: 文件类型（可选，如果不提供则从文件名推断）
        
        Returns:
            解析后的文本内容
        
        Raises:
            ValueError: 不支持的文件类型或解析失败
            RuntimeError: 缺少必要的依赖库
        """
        if not file_type:
            file_type = Path(file_path).suffix.lower().lstrip('.')
        
        parsers = {
            'pdf': cls.parse_pdf,
            'docx': cls.parse_docx,
            'xlsx': cls.parse_xlsx,
            'xls': cls.parse_xlsx,  # Excel旧格式也尝试用openpyxl
            'pptx': cls.parse_pptx,
            'txt': lambda p: Path(p).read_text(encoding='utf-8'),
            'md': lambda p: Path(p).read_text(encoding='utf-8'),
        }
        
        parser = parsers.get(file_type)
        if not parser:
            raise ValueError(
                f"不支持的文件类型: {file_type}。"
                f"支持的格式: {', '.join(parsers.keys())}"
            )
        
        try:
            logger.info(f"开始解析文件: {file_path} (类型: {file_type})")
            content = parser(file_path)
            logger.info(f"文件解析成功，提取 {len(content)} 字符")
            return content
        except RuntimeError as e:
            # 依赖缺失错误
            logger.error(f"依赖缺失: {e}")
            raise
        except Exception as e:
            logger.error(f"文件解析失败: {e}", exc_info=True)
            raise ValueError(f"文件解析失败: {str(e)}")

    @staticmethod
    def split_markdown_by_headings(markdown_text: str) -> List[Dict[str, Any]]:
        """
        按 Markdown 标题切分章节，返回带元数据的 section。

        每个元素结构：
        {
            "content": "章节文本",
            "meta": {
                "heading": "标题",
                "heading_level": 2,
                "heading_path": "父标题 > 子标题"
            }
        }
        """
        if not markdown_text or not markdown_text.strip():
            return []

        heading_pattern = re.compile(r'^(#{1,6})\s+(.*?)\s*$')
        lines = markdown_text.splitlines()

        sections: List[Dict[str, Any]] = []
        current_heading = "文档开始"
        current_level = 0
        current_lines: List[str] = []
        heading_stack: List[tuple[int, str]] = []
        in_code_fence = False

        def flush_section() -> None:
            nonlocal current_lines
            content = "\n".join(current_lines).strip()
            if not content:
                return

            heading_path = " > ".join([title for _, title in heading_stack]) if heading_stack else "文档开始"
            sections.append({
                "content": content,
                "meta": {
                    "heading": current_heading,
                    "heading_level": current_level,
                    "heading_path": heading_path,
                }
            })

        for line in lines:
            stripped = line.strip()

            # 处理代码块，避免把代码中的 # 误判为标题
            if stripped.startswith("```") or stripped.startswith("~~~"):
                in_code_fence = not in_code_fence
                current_lines.append(line)
                continue

            if not in_code_fence:
                match = heading_pattern.match(line)
                if match:
                    flush_section()

                    level = len(match.group(1))
                    title = match.group(2).strip() or "未命名标题"

                    while heading_stack and heading_stack[-1][0] >= level:
                        heading_stack.pop()
                    heading_stack.append((level, title))

                    current_heading = title
                    current_level = level
                    current_lines = []
                    continue

            current_lines.append(line)

        flush_section()
        return sections
